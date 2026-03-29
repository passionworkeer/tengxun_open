from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "reports" / "defense_deck_20260329.pptx"
IMG = ROOT / "img" / "final_delivery"


COLORS = {
    "navy": RGBColor(15, 23, 42),
    "slate": RGBColor(71, 85, 105),
    "muted": RGBColor(100, 116, 139),
    "sky": RGBColor(14, 165, 233),
    "mint": RGBColor(16, 185, 129),
    "amber": RGBColor(245, 158, 11),
    "rose": RGBColor(239, 68, 68),
    "ice": RGBColor(241, 245, 249),
    "white": RGBColor(255, 255, 255),
    "ink": RGBColor(15, 23, 42),
}


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    color: RGBColor | None = None,
    bold: bool = False,
    align: PP_ALIGN | None = None,
    font: str = "Aptos",
    margin: float = 0.02,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    p = box.text_frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color or COLORS["ink"]


def add_bullets(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    color: RGBColor | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.size = Pt(size)
        p.font.name = "Aptos"
        p.font.color.rgb = color or COLORS["ink"]


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    title: str,
    body: str,
    accent: RGBColor,
    fill: RGBColor = None,
) -> None:
    fill = fill or COLORS["white"]
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    add_text(slide, title, x + 0.18, y + 0.12, w - 0.36, 0.32, size=18, bold=True)
    add_text(
        slide,
        body,
        x + 0.18,
        y + 0.5,
        w - 0.36,
        h - 0.62,
        size=14,
        color=COLORS["slate"],
    )


def add_stat_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    label: str,
    value: str,
    accent: RGBColor,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = accent
    shape.line.width = Pt(2)
    add_text(
        slide,
        value,
        x + 0.15,
        y + 0.1,
        w - 0.3,
        0.48,
        size=22,
        color=accent,
        bold=True,
    )
    add_text(
        slide,
        label,
        x + 0.15,
        y + 0.58,
        w - 0.3,
        h - 0.68,
        size=13,
        color=COLORS["slate"],
    )


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.55, 0.28, 8.0, 0.55, size=28, bold=True)
    if subtitle:
        add_text(
            slide,
            subtitle,
            0.58,
            0.82,
            8.3,
            0.28,
            size=11,
            color=COLORS["muted"],
        )


def add_footer(slide, text: str) -> None:
    add_text(
        slide,
        text,
        0.55,
        5.1,
        8.8,
        0.18,
        size=9,
        color=COLORS["muted"],
    )


def add_picture(slide, rel_path: str, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(ROOT / rel_path), Inches(x), Inches(y), Inches(w), Inches(h))


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["navy"])
    add_text(
        slide,
        "Code Analysis Optimization with PE / RAG / FT",
        0.7,
        0.95,
        8.8,
        0.9,
        size=28,
        color=COLORS["white"],
        bold=True,
    )
    add_text(
        slide,
        "Celery cross-file dependency analysis, strict layered scoring, and ablation-driven strategy selection",
        0.75,
        2.0,
        8.4,
        0.6,
        size=15,
        color=COLORS["ice"],
    )
    add_stat_card(slide, 0.8, 3.0, 1.7, 1.0, label="Eval set", value="54", accent=COLORS["sky"])
    add_stat_card(slide, 2.7, 3.0, 1.7, 1.0, label="Few-shot", value="20", accent=COLORS["mint"])
    add_stat_card(slide, 4.6, 3.0, 1.7, 1.0, label="FT data", value="500", accent=COLORS["amber"])
    add_stat_card(slide, 6.5, 3.0, 2.2, 1.0, label="Best strict GPT PE", value="0.4757", accent=COLORS["rose"])
    add_text(
        slide,
        "Best business-model route: GPT-5.4 + postprocess_targeted",
        0.8,
        4.45,
        5.8,
        0.35,
        size=16,
        color=COLORS["white"],
        bold=True,
    )
    add_text(
        slide,
        "Branch: codex/strict-pe-search | March 29, 2026",
        0.8,
        5.25,
        4.8,
        0.2,
        size=10,
        color=COLORS["ice"],
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["ice"])
    add_header(slide, "Dataset and Evaluation Design", "Real open-source Celery cases, not synthetic toy samples")
    add_stat_card(slide, 0.6, 1.2, 2.0, 1.0, label="Manual eval cases", value="54", accent=COLORS["sky"])
    add_stat_card(slide, 2.85, 1.2, 2.0, 1.0, label="Strict few-shot", value="20", accent=COLORS["mint"])
    add_stat_card(slide, 5.1, 1.2, 2.0, 1.0, label="FT rows", value="500", accent=COLORS["amber"])
    add_stat_card(slide, 7.35, 1.2, 2.7, 1.0, label="Difficulty split", value="15 / 19 / 20", accent=COLORS["rose"])
    add_card(
        slide,
        0.65,
        2.55,
        4.1,
        1.45,
        title="Failure Types",
        body="Type A 7 | Type B 9 | Type C 11 | Type D 11 | Type E 16\n\nThe eval set emphasizes dynamic resolution, re-export chains, finalize/proxy flows, and ambiguity-heavy cases.",
        accent=COLORS["sky"],
    )
    add_card(
        slide,
        4.95,
        2.55,
        4.1,
        1.45,
        title="Scoring",
        body="Primary metric: union F1.\nStrict add-on: active-layer macro F1 and mislayer rate.\n\nThis exposes cases where the model finds the FQN but puts it in the wrong layer.",
        accent=COLORS["mint"],
    )
    add_bullets(
        slide,
        [
            "direct / indirect / implicit are preserved in the schema",
            "All 54 official cases are manually labeled against real source code",
            "Strict assets are used for audit, replay, and defense hardening",
        ],
        9.35,
        2.55,
        3.2,
        1.9,
        size=14,
    )
    add_footer(slide, "Official eval file: data/eval_cases.json")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])
    add_header(slide, "Why This Task Is Hard", "Cross-file dependency analysis fails for different reasons than simple import matching")
    add_card(
        slide,
        0.7,
        1.25,
        3.8,
        1.8,
        title="Long call chains",
        body="The model must follow multi-hop routes from CLI or public entry points through wrappers, re-exports, and factories before reaching the final stable symbol.",
        accent=COLORS["sky"],
    )
    add_card(
        slide,
        4.75,
        1.25,
        3.8,
        1.8,
        title="Implicit runtime edges",
        body="finalize hooks, Proxy resolution, decorator registration, string-based symbol lookup, and loader callbacks create dependencies that plain static matching misses.",
        accent=COLORS["mint"],
    )
    add_card(
        slide,
        8.8,
        1.25,
        3.8,
        1.8,
        title="Layer placement",
        body="Finding an FQN is not enough. The model must also decide whether it is direct, indirect, or implicit without flattening everything into one bucket.",
        accent=COLORS["amber"],
    )
    add_picture(slide, "img/final_delivery/03_bottleneck_heatmap_20260328.png", 0.95, 3.35, 5.6, 3.1)
    add_bullets(
        slide,
        [
            "Long-context loss",
            "Implicit edge omission",
            "Cross-file association breaks",
            "High union but wrong layer assignment",
        ],
        7.15,
        3.5,
        4.8,
        2.2,
        size=18,
    )
    add_footer(slide, "Key bottleneck framing: 'find the FQN' and 'place the layer' are different problems")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["ice"])
    add_header(slide, "Baseline Models", "Commercial ceiling vs. open-source starting point")
    add_picture(slide, "img/final_delivery/01_model_baselines_20260328.png", 0.65, 1.15, 7.0, 4.25)
    add_card(
        slide,
        8.0,
        1.35,
        4.6,
        1.15,
        title="GPT-5.4",
        body="Commercial upper bound.\nAvg = 0.2815",
        accent=COLORS["sky"],
    )
    add_card(
        slide,
        8.0,
        2.7,
        4.6,
        1.15,
        title="GLM-5",
        body="Mostly limited by structured-output adaptation.\nAvg = 0.0666",
        accent=COLORS["amber"],
    )
    add_card(
        slide,
        8.0,
        4.05,
        4.6,
        1.15,
        title="Qwen3.5-9B",
        body="Strict baseline is very low before systematic enhancement.\nAvg = 0.0370",
        accent=COLORS["rose"],
    )
    add_footer(slide, "Formal baseline table appears in reports/DELIVERY_REPORT.md")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])
    add_header(slide, "Prompt Engineering Is The Strongest Single Lever", "Official 54-case progressive GPT-5.4 run")
    add_picture(slide, "img/final_delivery/02_pe_progression_20260328.png", 0.65, 1.15, 7.25, 4.4)
    add_bullets(
        slide,
        [
            "Baseline 0.2745 -> postprocess 0.6062",
            "CoT is positive on the official 54-case line",
            "Few-shot is the main jump before post-processing",
            "But union F1 alone still overestimates layer quality",
        ],
        8.15,
        1.35,
        4.45,
        2.4,
        size=16,
    )
    add_stat_card(slide, 8.2, 4.2, 1.35, 0.9, label="Baseline", value="0.2745", accent=COLORS["slate"])
    add_stat_card(slide, 9.75, 4.2, 1.35, 0.9, label="Few-shot", value="0.5733", accent=COLORS["mint"])
    add_stat_card(slide, 11.3, 4.2, 1.35, 0.9, label="Postprocess", value="0.6062", accent=COLORS["sky"])
    add_footer(slide, "PE result files: results/pe_eval_54_20260328/")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["navy"])
    add_text(slide, "Why Strict Metrics Matter", 0.7, 0.35, 5.5, 0.5, size=26, color=COLORS["white"], bold=True)
    add_text(
        slide,
        "Union F1 can look high even when the model places dependencies in the wrong layer.",
        0.75,
        0.9,
        6.4,
        0.3,
        size=13,
        color=COLORS["ice"],
    )
    add_stat_card(slide, 0.8, 1.45, 2.2, 1.0, label="Old strict-best union", value="0.6136", accent=COLORS["sky"])
    add_stat_card(slide, 3.2, 1.45, 2.2, 1.0, label="Old strict-best macro", value="0.4372", accent=COLORS["amber"])
    add_stat_card(slide, 5.6, 1.45, 2.2, 1.0, label="Old mislayer", value="0.2336", accent=COLORS["rose"])
    add_stat_card(slide, 8.4, 1.45, 2.2, 1.0, label="New strict-best union", value="0.6338", accent=COLORS["sky"])
    add_stat_card(slide, 10.8, 1.45, 2.2, 1.0, label="New strict-best macro", value="0.4757", accent=COLORS["mint"])
    add_card(
        slide,
        0.8,
        3.0,
        3.7,
        1.5,
        title="Metric 1: active-layer macro F1",
        body="Penalizes wrong layer placement.\nA method must improve the active dependency layers, not only the union set.",
        accent=COLORS["mint"],
        fill=RGBColor(30, 41, 59),
    )
    add_card(
        slide,
        4.8,
        3.0,
        3.7,
        1.5,
        title="Metric 2: mislayer rate",
        body="Measures how often the model hits an FQN but assigns it to the wrong layer. Lower is better.",
        accent=COLORS["rose"],
        fill=RGBColor(30, 41, 59),
    )
    add_card(
        slide,
        8.8,
        3.0,
        3.7,
        1.5,
        title="Main lesson",
        body="The best PE strategy is the one that keeps union high while also increasing macro and reducing mislayer.",
        accent=COLORS["amber"],
        fill=RGBColor(30, 41, 59),
    )
    add_footer(slide, "Strict search report: reports/strict_pe_search_20260329.md")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["ice"])
    add_header(slide, "Strict PE Search: What Worked and What Failed", "The winning change was targeted example routing, not stronger prompt wording")
    add_card(slide, 0.7, 1.25, 3.0, 1.45, title="Failed: layer_guard", body="Higher union but much worse strict macro and mislayer.\nThe model became more pattern-driven, not more correct.", accent=COLORS["rose"])
    add_card(slide, 3.95, 1.25, 3.0, 1.45, title="Failed: assistant few-shot", body="Better JSON-following did not translate into better layer placement.\nStrict metrics regressed sharply.", accent=COLORS["rose"])
    add_card(slide, 7.2, 1.25, 2.8, 1.45, title="Worked: fewshot_targeted", body="Injected failure-mode anchors for Type B / E / D cases.\nImproved macro without exploding mislayer.", accent=COLORS["mint"])
    add_card(slide, 10.25, 1.25, 2.35, 1.45, title="Best: postprocess_targeted", body="Targeted selection plus layer-preserving parsing produced the new strict-best line.", accent=COLORS["sky"])
    add_bullets(
        slide,
        [
            "Type B anchors: shared_task / finalize / proxy / pending",
            "Type E anchors: symbol_by_name / loader / string import / backend",
            "Type D anchors: ambiguity-heavy registry or naming collisions",
            "Remaining few-shot slots are filled dynamically by relevance",
        ],
        0.9,
        3.15,
        5.7,
        1.8,
        size=15,
    )
    add_card(
        slide,
        6.9,
        3.1,
        5.2,
        1.8,
        title="Best strict PE numbers",
        body="postprocess_targeted\nunion 0.6338 | macro 0.4757 | mislayer 0.1620 | exact-layer 0.1296\n\nThis beats the prior strict-best 0.6136 / 0.4372 / 0.2336.",
        accent=COLORS["sky"],
    )
    add_footer(slide, "Result file: results/pe_targeted_full_20260329/pe_postprocess_targeted_strict.json")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])
    add_header(slide, "RAG Helps Hard Cases, Not Every Case", "RAG is a targeted fixer for long-chain and dynamic scenarios")
    add_picture(slide, "img/final_delivery/04_rag_retrieval_20260328.png", 0.6, 1.15, 5.9, 2.55)
    add_picture(slide, "img/final_delivery/05_rag_end_to_end_20260328.png", 0.6, 3.95, 5.9, 2.55)
    add_bullets(
        slide,
        [
            "Retrieval quality is measured separately: Recall@5 and MRR",
            "End-to-end gain is modest overall: 0.2783 -> 0.2940",
            "Hard cases improve strongly: 0.1980 -> 0.3372",
            "Conclusion: RAG should be enabled selectively, not by default",
        ],
        7.0,
        1.35,
        5.6,
        2.2,
        size=15,
    )
    add_card(
        slide,
        7.1,
        4.1,
        5.2,
        1.7,
        title="Operational recommendation",
        body="Turn on RAG for hard, long-chain, Type A, and Type E cases.\nDo not assume more context always means better code analysis quality.",
        accent=COLORS["amber"],
    )
    add_footer(slide, "Official RAG reports: reports/rag_pipeline.md and reports/rag_retrieval_eval_round4.md")

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["ice"])
    add_header(slide, "Open-Source Model Strategy: PE + FT Is The Core", "Fine-tuning alone is not enough")
    add_picture(slide, "img/final_delivery/06_qwen_strategies_20260328.png", 0.65, 1.2, 7.1, 4.5)
    add_stat_card(slide, 8.2, 1.35, 1.35, 0.9, label="FT only", value="0.0932", accent=COLORS["amber"])
    add_stat_card(slide, 9.75, 1.35, 1.35, 0.9, label="PE only", value="0.2246", accent=COLORS["sky"])
    add_stat_card(slide, 11.3, 1.35, 1.35, 0.9, label="PE + FT", value="0.3865", accent=COLORS["mint"])
    add_card(
        slide,
        8.1,
        2.65,
        4.7,
        1.45,
        title="Main takeaway",
        body="PE is the core gain source. FT learns domain patterns, but PE is what turns those patterns into stable, scorable JSON/FQN output.",
        accent=COLORS["mint"],
    )
    add_card(
        slide,
        8.1,
        4.35,
        4.7,
        1.15,
        title="Best open-source route",
        body="Highest strict-clean score: PE + RAG + FT = 0.5018\nLower-complexity strict route: PE + FT = 0.3865",
        accent=COLORS["sky"],
    )
    add_footer(slide, "Formal ablation matrix: reports/ablation_study.md")

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["navy"])
    add_text(slide, "Final Strategy Recommendation", 0.7, 0.45, 5.8, 0.45, size=26, color=COLORS["white"], bold=True)
    add_card(
        slide,
        0.8,
        1.35,
        3.8,
        2.0,
        title="Business model path",
        body="Use GPT-5.4 + postprocess_targeted.\nThis is the current strict-best PE route and the easiest line to defend methodologically.",
        accent=COLORS["sky"],
        fill=RGBColor(30, 41, 59),
    )
    add_card(
        slide,
        4.8,
        1.35,
        3.8,
        2.0,
        title="Open-source highest score",
        body="Use Qwen PE + RAG + FT as the current strict-clean highest-score route.\nIt is the strongest complete open-source path in the repo.",
        accent=COLORS["mint"],
        fill=RGBColor(30, 41, 59),
    )
    add_card(
        slide,
        8.8,
        1.35,
        3.8,
        2.0,
        title="Open-source default route",
        body="Use Qwen PE + FT as the lower-complexity strict-clean route.\nKeep historical 0.4315 only as an archived comparison point.",
        accent=COLORS["amber"],
        fill=RGBColor(30, 41, 59),
    )
    add_text(
        slide,
        "The key methodological conclusion is not 'stronger prompts win'. It is: failure-mode definition + targeted few-shot + layer-preserving postprocess wins.",
        0.9,
        4.25,
        11.5,
        0.55,
        size=19,
        color=COLORS["white"],
        bold=True,
    )
    add_text(
        slide,
        "Best strict GPT PE result: union 0.6338 | macro 0.4757 | mislayer 0.1620",
        0.95,
        5.05,
        7.3,
        0.25,
        size=13,
        color=COLORS["ice"],
    )

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])
    add_header(slide, "Appendix: Training Evidence and Boundary Conditions", "Useful if the defense moves toward FT quality or risk questions")
    add_picture(slide, "img/final_delivery/07_training_curve_20260328.png", 0.75, 1.2, 6.1, 3.8)
    add_bullets(
        slide,
        [
            "The training curve is stable and the final eval_loss is 0.4779",
            "Evidence for 'no overfitting' is moderate, not maximal, because the historical run did not keep a full step-wise eval curve",
            "The strict replay config now fixes this with eval_steps/save_steps=50 and best-checkpoint selection",
            "GLM official thinking mode was explored but not included in the formal matrix because both stream and non-stream smoke runs stalled on the first case",
            "Strict assets are the recommended path for replay and defense hardening",
        ],
        7.2,
        1.45,
        5.2,
        2.5,
        size=14,
    )
    add_card(
        slide,
        7.25,
        4.45,
        5.0,
        1.2,
        title="One-line boundary statement",
        body="RAG is a hard-case fixer, PE is the core gain source, and FT becomes effective mainly when combined with PE.",
        accent=COLORS["sky"],
    )
    add_footer(slide, "Supporting docs: reports/strict_replay_guide_20260329.md and reports/strict_data_audit_20260329.md")

    return prs


def main() -> int:
    prs = build_deck()
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.title = "Code Analysis Optimization Defense Deck"
    prs.core_properties.subject = "PE / RAG / FT on Celery code analysis"
    prs.core_properties.comments = "Generated from repository metrics and final delivery assets."
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
